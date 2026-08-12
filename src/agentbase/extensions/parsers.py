"""Document parsers for PDF, DOCX, HTML, Excel, and PPTX formats.

These parsers are registered with the global ``parser_registry`` at import
time. They use optional dependencies — if the library is not installed,
the parser will raise an ``ImportError`` with installation instructions
when ``parse()`` is called.

Dependencies:
- PDF:  ``pip install pymupdf``
- DOCX: ``pip install python-docx``
- HTML: ``pip install beautifulsoup4``
- Excel: ``pip install openpyxl``
- PPTX: ``pip install python-pptx``
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from agentbase.core.parsers import register_parser


@register_parser(".pdf", ".PDF", override=True)
class PdfParser:
    """PDF parser using pymupdf (fitz)."""

    extensions = [".pdf", ".PDF"]

    def parse(self, path: Path) -> str:
        try:
            import pymupdf  # type: ignore
        except ImportError:
            try:
                import fitz as pymupdf  # type: ignore
            except ImportError as exc:
                raise ImportError(
                    "PDF parsing requires pymupdf. Install with: pip install agentbase[rag]"
                ) from exc

        doc = pymupdf.open(str(path))
        pages: list[str] = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n\n".join(pages)


@register_parser(".docx", override=True)
class DocxParser:
    """DOCX parser using python-docx."""

    extensions = [".docx"]

    def parse(self, path: Path) -> str:
        try:
            import docx  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "DOCX parsing requires python-docx. Install with: pip install agentbase[rag]"
            ) from exc

        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)


@register_parser(".html", ".htm", override=True)
class HtmlParser:
    """HTML parser using beautifulsoup4."""

    extensions = [".html", ".htm"]

    def parse(self, path: Path) -> str:
        try:
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "HTML parsing requires beautifulsoup4. Install with: pip install beautifulsoup4"
            ) from exc

        html = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        # Remove script and style elements
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)


@register_parser(".xlsx", ".xls", override=True)
class ExcelParser:
    """Excel parser using openpyxl."""

    extensions = [".xlsx", ".xls"]

    def parse(self, path: Path) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "Excel parsing requires openpyxl. Install with: pip install openpyxl"
            ) from exc

        wb = load_workbook(str(path), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            lines.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append(" | ".join(cells))
        wb.close()
        return "\n".join(lines)


@register_parser(".pptx", ".ppt", override=True)
class PptxParser:
    """PPTX parser using python-pptx."""

    extensions = [".pptx", ".ppt"]

    def parse(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "PPTX parsing requires python-pptx. Install with: pip install python-pptx"
            ) from exc

        prs = Presentation(str(path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
            lines.append("")  # blank line between slides
        return "\n".join(lines)


# ---------------------------------------------------------------------------
# LLM-powered document parser (converts complex docs to structured Markdown)
# ---------------------------------------------------------------------------

@register_parser(".llm", override=True)
class LLMDocumentParser:
    """LLM-powered document parser.

    Uses a multimodal LLM (GPT-4o, Claude, etc.) to convert complex
    documents (scanned PDFs, images, complex layouts) into structured
    Markdown. Much higher quality than rule-based parsing, but costs
    API calls.

    Usage::

        # Register with a specific model
        from agentbase.extensions.parsers import LLMDocumentParser
        parser = LLMDocumentParser(model="gpt-4o", api_key="sk-...")

        # Or use via the parse() method with any file
        text = parser.parse(Path("complex_report.pdf"))

    The ``.llm`` extension is a virtual extension — call this parser
    directly for any file type that rule-based parsers handle poorly.
    """

    extensions = [".llm"]

    def __init__(
        self,
        *,
        model: str = "gpt-4o",
        api_key: str | None = None,
        base_url: str | None = None,
    ) -> None:
        self._model = model
        self._api_key = api_key
        self._base_url = base_url

    def _get_client(self):
        try:
            from openai import OpenAI
        except ImportError as exc:
            raise ImportError(
                "LLMDocumentParser requires the openai package. "
                "Install with: pip install openai"
            ) from exc
        import os
        kwargs: dict[str, Any] = {}
        key = self._api_key or os.environ.get("OPENAI_API_KEY") or os.environ.get("SILICONFLOW_API_KEY")
        if key:
            kwargs["api_key"] = key
        base = self._base_url or os.environ.get("OPENAI_BASE_URL")
        if base:
            kwargs["base_url"] = base
        return OpenAI(**kwargs)

    def parse(self, path: Path) -> str:
        import base64

        # Read file and encode as base64
        file_bytes = path.read_bytes()
        b64 = base64.b64encode(file_bytes).decode("utf-8")

        # Determine mime type from extension
        ext = path.suffix.lower()
        mime_map = {
            ".pdf": "application/pdf",
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }
        mime = mime_map.get(ext, "application/octet-stream")

        client = self._get_client()
        prompt = (
            "Convert this document into well-structured Markdown. "
            "Preserve headings, lists, tables, and code blocks. "
            "For images, describe what you see in alt text. "
            "Output only the Markdown content."
        )

        resp = client.chat.completions.create(
            model=self._model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "file" if mime == "application/pdf" else "image_url",
                            "image_url": {"url": f"data:{mime};base64,{b64}"} if mime != "application/pdf" else None,
                            "file": {"filename": path.name, "file_data": f"data:{mime};base64,{b64}"} if mime == "application/pdf" else None,
                        },
                    ],
                }
            ],
            temperature=0,
            max_tokens=4096,
        )
        return resp.choices[0].message.content or ""


# ---------------------------------------------------------------------------
# OCR parser for scanned PDFs and images
# ---------------------------------------------------------------------------

@register_parser(".ocr", override=True)
class OCRParser:
    """OCR parser for scanned documents and images.

    Uses ``pytesseract`` (Tesseract OCR) to extract text from images
    and scanned PDFs. Falls back to converting PDF pages to images
    via ``pdf2image`` first.

    Dependencies:
    - Tesseract OCR system binary (apt install tesseract-ocr)
    - ``pip install pytesseract pillow pdf2image``

    Usage::

        from agentbase.extensions.parsers import OCRParser
        parser = OCRParser()
        text = parser.parse(Path("scanned_document.pdf"))

    The ``.ocr`` extension is virtual — call this parser directly
    for any scanned/image file.
    """

    extensions = [".ocr"]

    def parse(self, path: Path) -> str:
        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise ImportError(
                "OCR parsing requires pytesseract and Pillow. "
                "Install with: pip install pytesseract pillow"
            ) from exc

        ext = path.suffix.lower()

        if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}:
            img = Image.open(path)
            return pytesseract.image_to_string(img, lang="chi_sim+eng")

        if ext == ".pdf":
            try:
                from pdf2image import convert_from_path
            except ImportError as exc:
                raise ImportError(
                    "PDF OCR requires pdf2image. "
                    "Install with: pip install pdf2image"
                ) from exc

            pages = convert_from_path(str(path))
            texts: list[str] = []
            for i, page in enumerate(pages, 1):
                text = pytesseract.image_to_string(page, lang="chi_sim+eng")
                texts.append(f"## Page {i}\n\n{text}")
            return "\n\n".join(texts)

        # For other files, try direct image OCR
        img = Image.open(path)
        return pytesseract.image_to_string(img, lang="chi_sim+eng")
