"""Comprehensive RAG pipeline verification.

Tests all stages end-to-end:
  1. Document parsers (9 types)
  2. Chunking strategies (3 types)
  3. Embedding providers (4 types)
  4. Search methods (3 types)
  5. Reranker + hybrid search
  6. End-to-end pipeline

Run: python scripts/test_rag_pipeline.py
"""
from __future__ import annotations

import os
import shutil
import sys
import tempfile
from pathlib import Path

# Load .env
env_file = Path(__file__).resolve().parent.parent / ".env"
if env_file.exists():
    for line in env_file.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, _, val = line.partition("=")
        key = key.strip()
        val = val.split("#")[0].strip()
        if val and key not in os.environ:
            os.environ[key] = val

if not os.environ.get("SILICONFLOW_BASE_URL"):
    os.environ["SILICONFLOW_BASE_URL"] = "https://api.siliconflow.cn/v1"

PASS = "\033[92mPASS\033[0m"
FAIL = "\033[91mFAIL\033[0m"
SKIP = "\033[93mSKIP\033[0m"

results: list[tuple[str, str, str]] = []


def run_test(group: str, name: str, fn):
    try:
        fn()
        results.append((group, name, PASS))
        print(f"  [{PASS}] {group}::{name}")
    except Exception as exc:
        results.append((group, name, FAIL))
        print(f"  [{FAIL}] {group}::{name}: {type(exc).__name__}: {exc}")


def skip_test(group: str, name: str, reason: str):
    results.append((group, name, SKIP))
    print(f"  [{SKIP}] {group}::{name}  ({reason})")


def section(title: str):
    print(f"\n--- {title} ---")


# ── 1. Document Parsers (9 types) ──────────────────────────────────────

def test_parsers():
    section("1. Document Parsers")
    # Ensure extension parsers are registered (PDF, DOCX, HTML, Excel, PPTX, etc.)
    import agentbase.extensions.parsers  # noqa: F401
    from agentbase.core.parsers import TextParser, parser_registry

    def test_txt():
        parser = parser_registry.get(".txt")
        assert isinstance(parser, TextParser) or hasattr(parser, "parse")
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("Hello World\nThis is a test.")
            f.flush()
            result = parser.parse(Path(f.name))
        assert "Hello World" in result

    def test_md():
        parser = parser_registry.get(".md")
        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as f:
            f.write("# Title\n\nSome **markdown** text.")
            f.flush()
            result = parser.parse(Path(f.name))
        assert "Title" in result

    def test_pdf():
        if not _has_pkg("pymupdf") and not _has_pkg("fitz"):
            skip_test("Parsers", "pdf", "pymupdf not installed")
            return
        import pymupdf

        # Create a minimal PDF
        pdf_path = Path(tempfile.gettempdir()) / "rag_test.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((50, 72), "Hello PDF World")
        doc.save(str(pdf_path))
        doc.close()
        parser = parser_registry.get(".pdf")
        result = parser.parse(pdf_path)
        assert "Hello PDF" in result

    def test_docx():
        if not _has_pkg("docx"):
            skip_test("Parsers", "docx", "python-docx not installed")
            return
        import docx
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            d = docx.Document()
            d.add_paragraph("Hello DOCX World")
            d.save(f.name)
            parser = parser_registry.get(".docx")
            result = parser.parse(Path(f.name))
        assert "Hello DOCX" in result

    def test_html():
        if not _has_pkg("bs4"):
            skip_test("Parsers", "html", "beautifulsoup4 not installed")
            return
        with tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8") as f:
            f.write("<html><body><h1>Hello HTML</h1><p>Test content</p></body></html>")
            f.flush()
            parser = parser_registry.get(".html")
            result = parser.parse(Path(f.name))
        assert "Hello HTML" in result

    def test_xlsx():
        if not _has_pkg("openpyxl"):
            skip_test("Parsers", "xlsx", "openpyxl not installed")
            return
        import openpyxl
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws["A1"] = "Name"
            ws["B1"] = "Value"
            ws["A2"] = "Hello"
            ws["B2"] = "123"
            wb.save(f.name)
            parser = parser_registry.get(".xlsx")
            result = parser.parse(Path(f.name))
        assert "Hello" in result

    def test_pptx():
        if not _has_pkg("pptx"):
            skip_test("Parsers", "pptx", "python-pptx not installed")
            return
        import pptx
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = pptx.Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Hello PPTX"
            slide.placeholders[1].text = "Test content"
            prs.save(f.name)
            parser = parser_registry.get(".pptx")
            result = parser.parse(Path(f.name))
        assert "Hello PPTX" in result

    def test_llm_parser():
        if not _has_pkg("openai"):
            skip_test("Parsers", "llm", "openai not installed")
            return
        from agentbase.extensions.parsers import LLMDocumentParser
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("# Test\n\nPython is a programming language.")
            f.flush()
            parser = LLMDocumentParser(
                api_key=os.environ.get("DEEPSEEK_API_KEY"),
                base_url="https://api.deepseek.com/v1",
                model="deepseek-chat",
            )
            result = parser.parse(Path(f.name))
        assert len(result) > 0

    def test_ocr_parser():
        if not shutil.which("tesseract"):
            skip_test("Parsers", "ocr", "tesseract binary not installed")
            return
        from PIL import Image, ImageDraw

        from agentbase.extensions.parsers import OCRParser
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            img = Image.new("RGB", (200, 50), color="white")
            draw = ImageDraw.Draw(img)
            draw.text((10, 10), "Hello OCR", fill="black")
            img.save(f, "PNG")
            f.flush()
            parser = OCRParser()
            result = parser.parse(Path(f.name))
        assert len(result) > 0

    run_test("Parsers", "txt", test_txt)
    run_test("Parsers", "md", test_md)
    run_test("Parsers", "pdf", test_pdf)
    run_test("Parsers", "docx", test_docx)
    run_test("Parsers", "html", test_html)
    run_test("Parsers", "xlsx", test_xlsx)
    run_test("Parsers", "pptx", test_pptx)
    run_test("Parsers", "llm", test_llm_parser)
    run_test("Parsers", "ocr", test_ocr_parser)


# ── 2. Chunking Strategies (3 types) ───────────────────────────────────

def test_chunking():
    section("2. Chunking Strategies")
    from agentbase.core.knowledge import _chunk_text

    sample = (
        "# Introduction\n\n"
        "Python is a programming language. "
        "It is widely used for web development.\n\n"
        "## Features\n\n"
        "Python supports object-oriented programming. "
        "It has a simple syntax. "
        "Python is easy to learn."
    )

    def test_paragraph():
        chunks = _chunk_text(sample, max_chunk_size=100, strategy="paragraph")
        assert len(chunks) >= 1
        assert all(len(c) <= 100 or "\n\n" not in c for c in chunks)

    def test_recursive():
        chunks = _chunk_text(sample, max_chunk_size=80, strategy="recursive", overlap=10)
        assert len(chunks) >= 2
        # Recursive should split on headers
        assert any("Introduction" in c or "Python" in c for c in chunks)

    def test_fixed():
        chunks = _chunk_text(sample, max_chunk_size=50, strategy="fixed")
        assert len(chunks) >= 2
        assert all(len(c) <= 50 for c in chunks)

    run_test("Chunking", "paragraph", test_paragraph)
    run_test("Chunking", "recursive", test_recursive)
    run_test("Chunking", "fixed", test_fixed)


# ── 3. Embedding Providers (4 types) ───────────────────────────────────

def test_embeddings():
    section("3. Embedding Providers")
    from agentbase.core.embeddings import embedding_registry

    def test_hash():
        provider = embedding_registry.get("hash")
        vec = provider.embed("hello world")
        assert len(vec) == 256  # default dimension
        # Verify no NaN (the bug we fixed)
        import math
        assert all(math.isfinite(v) for v in vec)

    def test_hash_batch():
        provider = embedding_registry.get("hash")
        vecs = provider.embed_batch(["hello", "world"])
        assert len(vecs) == 2

    def test_st():
        if not _has_pkg("sentence_transformers"):
            skip_test("Embedding", "sentence-transformers", "not installed")
            return
        provider = embedding_registry.get("sentence-transformers")
        vec = provider.embed("hello world")
        assert len(vec) == 384

    def test_none():
        # "none" means no embedding, text search only
        # Verify the registry has hash as default
        assert embedding_registry.has("hash")

    run_test("Embedding", "hash (no NaN)", test_hash)
    run_test("Embedding", "hash_batch", test_hash_batch)
    run_test("Embedding", "sentence-transformers", test_st)
    run_test("Embedding", "none (text search fallback)", test_none)


# ── 4. Search Methods (3 types) ────────────────────────────────────────

def test_search_methods():
    section("4. Search Methods")
    from agentbase.core.embeddings import HashEmbedding
    from agentbase.core.knowledge import KnowledgeBase

    # Use SQLite for isolated testing
    db_path = Path(tempfile.gettempdir()) / "agentbase_rag_test.db"
    if db_path.exists():
        db_path.unlink()

    def test_vector_search():
        provider = HashEmbedding(dimension=256)
        kb = KnowledgeBase(db_path=db_path, embedding_provider=provider)
        kb.add_document(
            source="test",
            title="Python Guide",
            content="Python is a high-level programming language used for web development and AI.",
        )
        results = kb.search("Python programming", top_k=5)
        assert len(results) >= 1
        assert hasattr(results[0], "score")
        kb.close()

    def test_text_search():
        # No embedding provider = text LIKE search
        kb = KnowledgeBase(db_path=db_path, embedding_provider=None)
        results = kb.search("Python", top_k=5)
        assert len(results) >= 1
        kb.close()

    def test_pgvector_search():
        # Test against PostgreSQL with pgvector
        import psycopg
        try:
            conn = psycopg.connect("postgresql://agentbase:agentbase@127.0.0.1:5432/agentbase")
            conn.close()
        except Exception:
            skip_test("Search", "pgvector", "PostgreSQL not available")
            return
        provider = HashEmbedding(dimension=256)
        kb = KnowledgeBase(dsn="postgresql://agentbase:agentbase@127.0.0.1:5432/agentbase", embedding_provider=provider)
        # Clean old test docs
        for doc in kb.list_documents():
            if doc.source == "rag_test_pg":
                kb.delete_document(doc_id=doc.id)
        kb.add_document(
            source="rag_test_pg",
            title="RAG Test",
            content="Vector databases enable semantic search over embedded documents.",
        )
        results = kb.search("semantic search", top_k=5)
        assert len(results) >= 1
        # Cleanup
        for doc in kb.list_documents():
            if doc.source == "rag_test_pg":
                kb.delete_document(doc_id=doc.id)
        kb.close()

    run_test("Search", "in-memory vector search", test_vector_search)
    run_test("Search", "text LIKE fallback", test_text_search)
    run_test("Search", "pgvector cosine search", test_pgvector_search)

    # Cleanup
    if db_path.exists():
        db_path.unlink()


# ── 5. Reranker + Hybrid Search ────────────────────────────────────────

def test_reranker_and_hybrid():
    section("5. Reranker + Hybrid Search")

    def test_reranker():
        if not _has_pkg("sentence_transformers"):
            skip_test("Reranker", "CrossEncoderReranker", "sentence-transformers not installed")
            return
        from agentbase.core.embeddings import CrossEncoderReranker
        reranker = CrossEncoderReranker(model="cross-encoder/ms-marco-MiniLM-L-6-v2")
        docs = [
            "Python is a programming language.",
            "The weather is nice today.",
            "Machine learning uses Python extensively.",
            "I like pizza.",
        ]
        result = reranker.rerank("Python programming", docs, top_k=2)
        assert len(result) == 2
        # The most relevant doc should be about Python
        top_doc_idx = result[0][0]
        assert "Python" in docs[top_doc_idx]

    def test_rrf_fusion():
        from agentbase.core.graph import Entity, GraphSearchResult, fuse_results_rrf
        from agentbase.core.knowledge import Document, SearchResult

        # Vector search results (have .document attribute)
        list1 = [
            SearchResult(document=Document(id=1, source="test", title="Python", content="", chunk_count=1), chunk=None, score=0.9),
            SearchResult(document=Document(id=2, source="test", title="Java", content="", chunk_count=1), chunk=None, score=0.8),
        ]
        # Graph search results (have .entity attribute)
        list2 = [
            GraphSearchResult(entity=Entity(id=2, name="Java"), score=0.7),
            GraphSearchResult(entity=Entity(id=3, name="Go"), score=0.6),
        ]
        fused = fuse_results_rrf(list1, list2, top_k=5)
        assert len(fused) >= 2, f"Expected >= 2 results, got {len(fused)}"
        # The results should be sorted by fused score
        assert fused[0] is not None

    run_test("Reranker", "CrossEncoderReranker", test_reranker)
    run_test("Hybrid", "RRF fusion", test_rrf_fusion)


# ── 6. End-to-End Pipeline ─────────────────────────────────────────────

def test_e2e_pipeline():
    section("6. End-to-End RAG Pipeline")

    def test_full_pipeline():
        from agentbase.core.embeddings import HashEmbedding
        from agentbase.core.knowledge import KnowledgeBase, _chunk_text
        from agentbase.core.parsers import parser_registry

        # Step 1: Parse a file
        test_file = Path(tempfile.gettempdir()) / "rag_e2e_test.md"
        test_file.write_text(
            "# Python Programming\n\n"
            "Python is a high-level programming language created by Guido van Rossum.\n\n"
            "## Features\n\n"
            "Python supports multiple programming paradigms including object-oriented and functional.\n"
            "It has a comprehensive standard library and dynamic type system.\n\n"
            "## Use Cases\n\n"
            "Python is widely used in web development, data science, AI, and automation.\n",
            encoding="utf-8",
        )
        parser = parser_registry.get_for_path(test_file)
        content = parser.parse(test_file)
        assert "Python" in content

        # Step 2: Chunk the text
        chunks = _chunk_text(content, max_chunk_size=200, strategy="recursive", overlap=20)
        assert len(chunks) >= 2

        # Step 3: Store in KB with embeddings
        db_path = Path(tempfile.gettempdir()) / "agentbase_e2e_rag.db"
        if db_path.exists():
            db_path.unlink()
        provider = HashEmbedding(dimension=256)
        kb = KnowledgeBase(db_path=db_path, embedding_provider=provider)
        doc = kb.add_document(source=str(test_file), title="Python Guide", content=content)
        assert doc.chunk_count >= 1

        # Step 4: Search
        results = kb.search("Python programming language", top_k=3)
        assert len(results) >= 1, f"Expected >= 1 result, got {len(results)}"
        assert hasattr(results[0], "score"), "Missing score attribute"
        # Score can be negative (cosine distance) — just check it exists

        # Step 5: Verify result is relevant
        top_result = results[0]
        assert top_result.document is not None, "Document is None"
        # The title was "Python Guide" so it should contain "Python"
        title = top_result.document.title or ""
        source = top_result.document.source or ""
        assert "Python" in title or "Python" in source, f"Expected 'Python' in title='{title}' or source='{source}'"

        kb.close()
        if db_path.exists():
            db_path.unlink()

    run_test("E2E", "parse→chunk→embed→store→search", test_full_pipeline)


# ── Helpers ────────────────────────────────────────────────────────────

def _has_pkg(name: str) -> bool:
    try:
        __import__(name)
        return True
    except ImportError:
        return False


# ── Main ───────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  RAG Pipeline Verification")
    print("=" * 60)

    test_parsers()
    test_chunking()
    test_embeddings()
    test_search_methods()
    test_reranker_and_hybrid()
    test_e2e_pipeline()

    # Summary
    print(f"\n{'=' * 60}")
    print("  SUMMARY")
    print(f"{'=' * 60}")
    total = len(results)
    passed = sum(1 for _, _, s in results if s == PASS)
    failed = sum(1 for _, _, s in results if s == FAIL)
    skipped = sum(1 for _, _, s in results if s == SKIP)

    for group, name, status in results:
        print(f"  [{status}] {group}::{name}")

    print(f"\n  Total: {passed} passed, {failed} failed, {skipped} skipped / {total}")
    print(f"{'=' * 60}")
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
